from __future__ import annotations

from datetime import UTC, datetime
from io import BytesIO
from pathlib import Path
from typing import Any
import zipfile
from xml.sax.saxutils import escape as xml_escape

from orchid_ai.core.tool import OrchidTool, OrchidToolInput, OrchidToolOutput

from .write_file import _ensure_extension, _write_content

_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def _stringify_content(content: Any) -> str:
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        return "\n".join(_stringify_content(item) for item in content)
    if isinstance(content, dict):
        parts = [f"{key}: {_stringify_content(value)}" for key, value in content.items()]
        return "\n".join(parts)
    return str(content)


def _coerce_slides(slides: list[Any], title: str, filename: str) -> list[dict[str, str]]:
    deck_title = title or Path(filename).stem or "Presentation"
    normalised: list[dict[str, str]] = [{"title": deck_title, "content": ""}]
    for index, slide in enumerate(slides, start=1):
        if isinstance(slide, dict):
            slide_title = str(slide.get("title") or f"Slide {index}")
            slide_content = _stringify_content(slide.get("content") or slide.get("body") or "")
        else:
            slide_title = f"Slide {index}"
            slide_content = _stringify_content(slide)
        normalised.append({"title": slide_title, "content": slide_content})
    return normalised


def _build_with_python_pptx(slides: list[dict[str, str]]) -> bytes:
    from pptx import Presentation

    presentation = Presentation()

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = slides[0]["title"]
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = slides[0]["content"]

    for slide in slides[1:]:
        ppt_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        ppt_slide.shapes.title.text = slide["title"]
        body_placeholder = ppt_slide.placeholders[1]
        text_frame = body_placeholder.text_frame
        lines = [line for line in slide["content"].splitlines() if line.strip()]
        text_frame.text = lines[0] if lines else ""
        for line in lines[1:]:
            text_frame.add_paragraph().text = line

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _empty_group_properties() -> str:
    return (
        "<p:nvGrpSpPr>"
        '<p:cNvPr id="1" name=""/>'
        "<p:cNvGrpSpPr/>"
        "<p:nvPr/>"
        "</p:nvGrpSpPr>"
        "<p:grpSpPr>"
        "<a:xfrm>"
        '<a:off x="0" y="0"/>'
        '<a:ext cx="0" cy="0"/>'
        '<a:chOff x="0" y="0"/>'
        '<a:chExt cx="0" cy="0"/>'
        "</a:xfrm>"
        "</p:grpSpPr>"
    )


def _shape_xml(shape_id: int, name: str, text: str, *, x: int, y: int, cx: int, cy: int, bold: bool = False) -> str:
    paragraphs = []
    for line in [line for line in text.splitlines() if line.strip()] or [""]:
        run_props = ' b="1"' if bold else ""
        paragraphs.append(
            "<a:p>"
            f'<a:r><a:rPr lang="en-US"{run_props}/><a:t>{xml_escape(line)}</a:t></a:r>'
            '<a:endParaRPr lang="en-US"/>'
            "</a:p>"
        )
    return (
        "<p:sp>"
        "<p:nvSpPr>"
        f'<p:cNvPr id="{shape_id}" name="{xml_escape(name)}"/>'
        "<p:cNvSpPr/>"
        "<p:nvPr/>"
        "</p:nvSpPr>"
        "<p:spPr>"
        "<a:xfrm>"
        f'<a:off x="{x}" y="{y}"/>'
        f'<a:ext cx="{cx}" cy="{cy}"/>'
        "</a:xfrm>"
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        "</p:spPr>"
        '<p:txBody><a:bodyPr wrap="square"/><a:lstStyle/>'
        + "".join(paragraphs)
        + "</p:txBody>"
        "</p:sp>"
    )


def _slide_xml(title: str, content: str) -> str:
    title_shape = _shape_xml(2, "Title", title, x=457200, y=274638, cx=8229600, cy=914400, bold=True)
    body_shape = _shape_xml(3, "Content", content, x=457200, y=1371600, cx=8229600, cy=4572000)
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<p:sld xmlns:a="{_A_NS}" xmlns:r="{_R_NS}" xmlns:p="{_P_NS}">'
        "<p:cSld>"
        "<p:spTree>"
        + _empty_group_properties()
        + title_shape
        + body_shape
        + "</p:spTree>"
        "</p:cSld>"
        "<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>"
        "</p:sld>"
    )


def _slide_layout_xml() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<p:sldLayout xmlns:a="{_A_NS}" xmlns:r="{_R_NS}" xmlns:p="{_P_NS}" type="blank" preserve="1">'
        "<p:cSld name=\"Blank\">"
        "<p:spTree>"
        + _empty_group_properties()
        + "</p:spTree>"
        "</p:cSld>"
        "<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>"
        "</p:sldLayout>"
    )


def _slide_master_xml() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<p:sldMaster xmlns:a="{_A_NS}" xmlns:r="{_R_NS}" xmlns:p="{_P_NS}">'
        "<p:cSld name=\"Master\">"
        "<p:spTree>"
        + _empty_group_properties()
        + "</p:spTree>"
        "</p:cSld>"
        '<p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" '
        'accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>'
        "<p:sldLayoutIdLst><p:sldLayoutId id=\"2147483649\" r:id=\"rId1\"/></p:sldLayoutIdLst>"
        "<p:txStyles>"
        "<p:titleStyle/>"
        "<p:bodyStyle/>"
        "<p:otherStyle/>"
        "</p:txStyles>"
        "</p:sldMaster>"
    )


def _theme_xml() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<a:theme xmlns:a="{_A_NS}" name="Orchid">'
        "<a:themeElements>"
        "<a:clrScheme name=\"Orchid\">"
        "<a:dk1><a:srgbClr val=\"000000\"/></a:dk1>"
        "<a:lt1><a:srgbClr val=\"FFFFFF\"/></a:lt1>"
        "<a:dk2><a:srgbClr val=\"1F2937\"/></a:dk2>"
        "<a:lt2><a:srgbClr val=\"F9FAFB\"/></a:lt2>"
        "<a:accent1><a:srgbClr val=\"2563EB\"/></a:accent1>"
        "<a:accent2><a:srgbClr val=\"059669\"/></a:accent2>"
        "<a:accent3><a:srgbClr val=\"D97706\"/></a:accent3>"
        "<a:accent4><a:srgbClr val=\"DC2626\"/></a:accent4>"
        "<a:accent5><a:srgbClr val=\"7C3AED\"/></a:accent5>"
        "<a:accent6><a:srgbClr val=\"0F766E\"/></a:accent6>"
        "<a:hlink><a:srgbClr val=\"2563EB\"/></a:hlink>"
        "<a:folHlink><a:srgbClr val=\"7C3AED\"/></a:folHlink>"
        "</a:clrScheme>"
        "<a:fontScheme name=\"Orchid\">"
        '<a:majorFont><a:latin typeface="Aptos"/></a:majorFont>'
        '<a:minorFont><a:latin typeface="Aptos"/></a:minorFont>'
        "</a:fontScheme>"
        "<a:fmtScheme name=\"Orchid\">"
        "<a:fillStyleLst><a:solidFill><a:schemeClr val=\"accent1\"/></a:solidFill></a:fillStyleLst>"
        "<a:lnStyleLst><a:ln w=\"9525\"><a:solidFill><a:schemeClr val=\"accent1\"/></a:solidFill></a:ln></a:lnStyleLst>"
        "<a:effectStyleLst><a:effectStyle/></a:effectStyleLst>"
        '<a:bgFillStyleLst><a:solidFill><a:schemeClr val="lt1"/></a:solidFill></a:bgFillStyleLst>'
        "</a:fmtScheme>"
        "</a:themeElements>"
        "<a:objectDefaults/>"
        "<a:extraClrSchemeLst/>"
        "</a:theme>"
    )


def _content_types_xml(slide_count: int) -> str:
    slide_overrides = "".join(
        f'<Override PartName="/ppt/slides/slide{index}.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        for index in range(1, slide_count + 1)
    )
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>'
        '<Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>'
        '<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
        '<Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>'
        '<Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>'
        '<Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>'
        + slide_overrides
        + "</Types>"
    )


def _root_relationships_xml() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<Relationships xmlns="{_REL_NS}">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>'
        '<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>'
        '<Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>'
        "</Relationships>"
    )


def _presentation_xml(slide_count: int) -> str:
    slide_ids = "".join(
        f'<p:sldId id="{255 + index}" r:id="rId{index + 1}"/>'
        for index in range(1, slide_count + 1)
    )
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<p:presentation xmlns:a="{_A_NS}" xmlns:r="{_R_NS}" xmlns:p="{_P_NS}">'
        '<p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>'
        f"<p:sldIdLst>{slide_ids}</p:sldIdLst>"
        '<p:sldSz cx="9144000" cy="6858000"/>'
        '<p:notesSz cx="6858000" cy="9144000"/>'
        "<p:defaultTextStyle/>"
        "</p:presentation>"
    )


def _presentation_relationships_xml(slide_count: int) -> str:
    relationships = [
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>'
    ]
    for index in range(1, slide_count + 1):
        relationships.append(
            f'<Relationship Id="rId{index + 1}" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
            f'Target="slides/slide{index}.xml"/>'
        )
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<Relationships xmlns="{_REL_NS}">'
        + "".join(relationships)
        + "</Relationships>"
    )


def _slide_relationships_xml() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<Relationships xmlns="{_REL_NS}">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>'
        "</Relationships>"
    )


def _slide_master_relationships_xml() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<Relationships xmlns="{_REL_NS}">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>'
        '<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>'
        "</Relationships>"
    )


def _slide_layout_relationships_xml() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<Relationships xmlns="{_REL_NS}">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/>'
        "</Relationships>"
    )


def _app_properties_xml(slide_count: int) -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" '
        'xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">'
        "<Application>Orchid</Application>"
        "<PresentationFormat>Custom</PresentationFormat>"
        f"<Slides>{slide_count}</Slides>"
        "<Notes>0</Notes>"
        "<HiddenSlides>0</HiddenSlides>"
        "</Properties>"
    )


def _core_properties_xml(title: str) -> str:
    timestamp = datetime.now(UTC).replace(microsecond=0).isoformat().replace("+00:00", "Z")
    safe_title = xml_escape(title)
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" '
        'xmlns:dc="http://purl.org/dc/elements/1.1/" '
        'xmlns:dcterms="http://purl.org/dc/terms/" '
        'xmlns:dcmitype="http://purl.org/dc/dcmitype/" '
        'xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">'
        f"<dc:title>{safe_title}</dc:title>"
        "<dc:creator>Orchid</dc:creator>"
        "<cp:lastModifiedBy>Orchid</cp:lastModifiedBy>"
        f'<dcterms:created xsi:type="dcterms:W3CDTF">{timestamp}</dcterms:created>'
        f'<dcterms:modified xsi:type="dcterms:W3CDTF">{timestamp}</dcterms:modified>'
        "</cp:coreProperties>"
    )


def _build_pptx_fallback(slides: list[dict[str, str]]) -> bytes:
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", _content_types_xml(len(slides)))
        archive.writestr("_rels/.rels", _root_relationships_xml())
        archive.writestr("docProps/app.xml", _app_properties_xml(len(slides)))
        archive.writestr("docProps/core.xml", _core_properties_xml(slides[0]["title"]))
        archive.writestr("ppt/presentation.xml", _presentation_xml(len(slides)))
        archive.writestr("ppt/_rels/presentation.xml.rels", _presentation_relationships_xml(len(slides)))
        archive.writestr("ppt/slideMasters/slideMaster1.xml", _slide_master_xml())
        archive.writestr("ppt/slideMasters/_rels/slideMaster1.xml.rels", _slide_master_relationships_xml())
        archive.writestr("ppt/slideLayouts/slideLayout1.xml", _slide_layout_xml())
        archive.writestr("ppt/slideLayouts/_rels/slideLayout1.xml.rels", _slide_layout_relationships_xml())
        archive.writestr("ppt/theme/theme1.xml", _theme_xml())
        for index, slide in enumerate(slides, start=1):
            archive.writestr(f"ppt/slides/slide{index}.xml", _slide_xml(slide["title"], slide["content"]))
            archive.writestr(f"ppt/slides/_rels/slide{index}.xml.rels", _slide_relationships_xml())
    return buffer.getvalue()


class GeneratePPTXTool(OrchidTool):
    name = "generate_pptx"
    description = "Generate a PPTX deck with a title slide and one slide per item"
    parameters_schema = {
        "type": "object",
        "properties": {
            "slides": {
                "type": "array",
                "description": "Slides to render, each with a title and content",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "content": {"type": "string"},
                    },
                },
            },
            "filename": {
                "type": "string",
                "description": "Output filename without extension",
            },
            "title": {
                "type": "string",
                "description": "Optional presentation title",
                "default": "",
            },
        },
        "required": ["slides", "filename"],
    }
    parallel_safe = False

    async def invoke(self, tool_input: OrchidToolInput) -> OrchidToolOutput:
        parameters = tool_input.parameters
        slides = _coerce_slides(
            slides=list(parameters.get("slides") or []),
            title=str(parameters.get("title", "") or ""),
            filename=str(parameters["filename"]),
        )

        try:
            pptx_bytes = _build_with_python_pptx(slides)
            generator = "python-pptx"
        except ModuleNotFoundError:
            pptx_bytes = _build_pptx_fallback(slides)
            generator = "fallback"

        filepath = str(_ensure_extension(str(parameters["filename"]), ".pptx"))
        path, size = _write_content(
            pptx_bytes,
            filepath=filepath,
            mode="binary",
            context=tool_input.context,
            content_sources=tool_input.content_sources,
            auth_context=tool_input.auth_context,
        )
        return OrchidToolOutput(
            result={"path": str(path), "size_bytes": size, "format": "pptx"},
            metadata={
                "generator": generator,
                "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            },
        )
